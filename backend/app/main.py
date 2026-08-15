import os
import uuid
import shutil
from fastapi import FastAPI, UploadFile, File, HTTPException, BackgroundTasks, Query, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, StreamingResponse
from pydantic import BaseModel
from typing import List, Dict, Any, Optional

from backend.app.config import settings
from backend.app.rag import rag_service, parse_document_ocr
from backend.app.generator.paper_generator import generate_full_paper, JOURNAL_TIERS
from backend.app.advisor.plagiarism_advisor import generate_reduction_advice
from backend.app import database as db
from backend.app.models import (
    PaperResponse, 
    PaperAnalysisResponse, 
    SectionRewriteRequest, 
    PaperRewriteRequest,
    SectionResponse
)
from backend.app.parser import parse_document
from backend.app.chunker.chunk_manager import get_paragraph_with_context
from backend.app.similarity import analyze_paper_similarity, save_section_embeddings
from contextlib import asynccontextmanager

from backend.app.rewrite import rewrite_text
from backend.app.quality import check_academic_quality
from backend.app.rebuild import rebuild_document
from backend.app.export.exporter import convert_docx_to_pdf
from backend.app.similarity.citation_resolver import resolve_and_format_citation
from backend.app.similarity.pipeline_worker import run_autonomous_pipeline



@asynccontextmanager
async def lifespan(app: FastAPI):
    db.init_db()
    yield

app = FastAPI(title=settings.PROJECT_NAME, lifespan=lifespan)

# Setup CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

@app.post("/api/v1/papers/upload", response_model=PaperResponse)
async def upload_paper(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    journal_format: str = Query("ieee")
):
    ext = os.path.splitext(file.filename)[1].lower()
    if ext not in [".docx", ".pdf"]:
        raise HTTPException(status_code=400, detail="Only DOCX and PDF files are supported.")
        
    paper_id = str(uuid.uuid4())
    temp_filename = f"{paper_id}{ext}"
    temp_path = os.path.join(settings.UPLOAD_DIR, temp_filename)
    
    with open(temp_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
        
    try:
        # Create paper record in DB with queued status
        db.create_paper(paper_id, file.filename, ext.strip("."))
        db.update_paper_status(paper_id, "queued")
        
        # Hand off process to the autonomous pipeline background worker
        background_tasks.add_task(
            run_autonomous_pipeline,
            paper_id=paper_id,
            temp_path=temp_path,
            journal_format=journal_format
        )
        
        return db.get_paper(paper_id)
    except Exception as e:
        # Clean up on failure
        if os.path.exists(temp_path):
            try:
                os.remove(temp_path)
            except:
                pass
        db.update_paper_status(paper_id, f"error: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Failed to queue document: {str(e)}")

@app.get("/api/v1/papers", response_model=List[PaperResponse])
def get_papers():
    return db.get_all_papers()

@app.get("/api/v1/papers/{paper_id}", response_model=PaperAnalysisResponse)
def get_paper_details(paper_id: str):
    paper = db.get_paper(paper_id)
    if not paper:
        raise HTTPException(status_code=404, detail="Paper not found.")
        
    sections = db.get_paper_sections(paper_id)
    references = db.get_paper_references(paper_id)
    
    return {
        "paper": paper,
        "sections": sections,
        "references": references
    }

@app.post("/api/v1/papers/{paper_id}/analyze")
def run_similarity_analysis(paper_id: str):
    paper = db.get_paper(paper_id)
    if not paper:
        raise HTTPException(status_code=404, detail="Paper not found.")
        
    sections = db.get_paper_sections(paper_id)
    if not sections:
        raise HTTPException(status_code=400, detail="Paper contains no text blocks to analyze.")
        
    # Analyze similarity against other papers
    result = analyze_paper_similarity(paper_id, sections)
    return {
        "overall_similarity": result["overall_similarity"],
        "flagged_count": result["flagged_count"]
    }

@app.post("/api/v1/sections/{section_id}/rewrite", response_model=SectionResponse)
def rewrite_section(section_id: str, request: SectionRewriteRequest):
    section = db.get_section(section_id)
    if not section:
        raise HTTPException(status_code=404, detail="Section not found.")
        
    # Get paragraph with context for LLM rewriter
    sections = db.get_paper_sections(section["paper_id"])
    target_idx = next(i for i, s in enumerate(sections) if s["id"] == section_id)
    context_info = get_paragraph_with_context(sections, target_idx)
    
    # Run LLM rewrite
    rewritten_text = rewrite_text(
        text=context_info["text"],
        context_before=context_info["context_before"],
        context_after=context_info["context_after"],
        target_similarity=settings.SIMILARITY_THRESHOLD
    )
    
    # Run quality validation
    quality_report = check_academic_quality(context_info["text"], rewritten_text)
    
    # Update layout metadata with warnings if any
    meta = section["layout_metadata"]
    meta["quality_warnings"] = quality_report["warnings"]
    
    # Save rewrite to DB (reset similarity score for re-checking)
    db.update_section_rewrite(section_id, rewritten_text, similarity_score=0.0, is_flagged=False)
    
    return db.get_section(section_id)

@app.post("/api/v1/papers/{paper_id}/rewrite-all")
def rewrite_flagged_sections(paper_id: str, background_tasks: BackgroundTasks):
    paper = db.get_paper(paper_id)
    if not paper:
        raise HTTPException(status_code=404, detail="Paper not found.")
        
    sections = db.get_paper_sections(paper_id)
    flagged_sections = [s for s in sections if s["is_flagged"]]
    
    if not flagged_sections:
        return {"message": "No flagged sections found to rewrite."}
        
    def process_bulk_rewrites():
        db.update_paper_status(paper_id, "rewriting")
        for idx, sec in enumerate(flagged_sections):
            try:
                target_idx = next(i for i, s in enumerate(sections) if s["id"] == sec["id"])
                context_info = get_paragraph_with_context(sections, target_idx)
                
                # Rewrite text
                rewritten_text = rewrite_text(
                    text=context_info["text"],
                    context_before=context_info["context_before"],
                    context_after=context_info["context_after"]
                )
                
                # Check quality
                quality_report = check_academic_quality(context_info["text"], rewritten_text)
                meta = sec["layout_metadata"]
                meta["quality_warnings"] = quality_report["warnings"]
                
                db.update_section_rewrite(sec["id"], rewritten_text, similarity_score=0.0, is_flagged=False)
            except Exception as e:
                print(f"Error bulk rewriting section {sec['id']}: {e}")
                
        # Run similarity recheck after all rewrites are complete
        updated_sections = db.get_paper_sections(paper_id)
        analyze_paper_similarity(paper_id, updated_sections)
        db.update_paper_status(paper_id, "ready")

    background_tasks.add_task(process_bulk_rewrites)
    return {"message": f"Started background rewrite of {len(flagged_sections)} flagged sections."}

@app.post("/api/v1/papers/{paper_id}/rebuild")
def rebuild_paper_docx(paper_id: str, request: PaperRewriteRequest):
    paper = db.get_paper(paper_id)
    if not paper:
        raise HTTPException(status_code=404, detail="Paper not found.")
        
    sections = db.get_paper_sections(paper_id)
    references = db.get_paper_references(paper_id)
    
    # Simple layout mapping
    layout_map = {
        "margins": {"top": 1.0, "bottom": 1.0, "left": 1.0, "right": 1.0}
    }
    
    output_filename = f"rebuilt_{paper_id}.docx"
    output_path = os.path.join(settings.OUTPUT_DIR, output_filename)
    
    try:
        rebuild_document(
            sections=sections,
            references=references,
            layout_map=layout_map,
            output_path=output_path,
            journal_format=request.journal_format,
            media_dir=os.path.join(settings.UPLOAD_DIR, "media")
        )
        return {"filename": output_filename, "format": "docx"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to rebuild document: {str(e)}")

@app.post("/api/v1/papers/{paper_id}/resolve-references")
async def resolve_references_endpoint(
    paper_id: str, 
    style: str = Query("numeric", enum=["numeric", "author_year"])
):
    paper = db.get_paper(paper_id)
    if not paper:
        raise HTTPException(status_code=404, detail="Paper not found.")
        
    references = db.get_paper_references(paper_id)
    if not references:
        return {"message": "No references found to resolve."}
        
    resolved_count = 0
    for idx, ref in enumerate(references):
        # Query Semantic Scholar to resolve metadata and format the citation
        res = await resolve_and_format_citation(
            raw_reference=ref["raw_reference"],
            style=style,
            citation_idx=idx + 1
        )
        
        # If resolved, update the database row with formatted text and DOI as citation key
        if res["resolved"]:
            cit_key = f"DOI: {res['doi']}" if res["doi"] else ref["citation_key"]
            db.update_reference(ref["id"], res["formatted_reference"], cit_key)
            resolved_count += 1
            
    return {
        "message": f"Successfully resolved and formatted {resolved_count} of {len(references)} references.",
        "resolved_count": resolved_count,
        "total_count": len(references)
    }


@app.get("/api/v1/papers/{paper_id}/download")
def download_rebuilt_file(paper_id: str, file_format: str = Query("docx", enum=["docx", "pdf"])):
    docx_filename = f"rebuilt_{paper_id}.docx"
    docx_path = os.path.join(settings.OUTPUT_DIR, docx_filename)
    
    if not os.path.exists(docx_path):
        raise HTTPException(status_code=404, detail="Rebuilt document not found. Please run rebuild first.")
        
    if file_format == "pdf":
        pdf_filename = f"rebuilt_{paper_id}.pdf"
        pdf_path = os.path.join(settings.OUTPUT_DIR, pdf_filename)
        
        # If PDF already generated, serve it
        if os.path.exists(pdf_path):
            return FileResponse(pdf_path, media_type="application/pdf", filename=f"rebuilt_paper.pdf")
            
        # Try to convert on the fly
        converted = convert_docx_to_pdf(docx_path, pdf_path)
        if converted and os.path.exists(pdf_path):
            return FileResponse(pdf_path, media_type="application/pdf", filename=f"rebuilt_paper.pdf")
        else:
            raise HTTPException(
                status_code=501, 
                detail="PDF conversion engine is not available on this server. Please download the DOCX file."
            )
            
    return FileResponse(docx_path, media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document", filename=f"rebuilt_paper.docx")


# --- RAG & OCR API Endpoints ---

class RAGQueryRequest(BaseModel):
    query: str
    paper_id: Optional[str] = None
    top_k: Optional[int] = 4

@app.post("/api/v1/rag/query")
def rag_query(request: RAGQueryRequest):
    if not request.query.strip():
        raise HTTPException(status_code=400, detail="Query string cannot be empty.")
    return rag_service.query_rag(query=request.query, paper_id=request.paper_id, top_k=request.top_k or 4)

@app.post("/api/v1/rag/stream")
def rag_stream_query(request: RAGQueryRequest):
    if not request.query.strip():
        raise HTTPException(status_code=400, detail="Query string cannot be empty.")
    return StreamingResponse(
        rag_service.stream_rag_response(query=request.query, paper_id=request.paper_id, top_k=request.top_k or 4),
        media_type="text/event-stream"
    )

@app.post("/api/v1/rag/ocr")
async def rag_ocr_upload(file: UploadFile = File(...)):
    ext = os.path.splitext(file.filename)[1].lower()
    if ext not in [".pdf", ".png", ".jpg", ".jpeg"]:
        raise HTTPException(status_code=400, detail="Unsupported file format for OCR.")
    
    temp_id = str(uuid.uuid4())
    temp_path = os.path.join(settings.UPLOAD_DIR, f"ocr_{temp_id}{ext}")
    
    with open(temp_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
        
    try:
        parsed_result = parse_document_ocr(temp_path)
        return parsed_result
    finally:
        if os.path.exists(temp_path):
            os.remove(temp_path)


# --- Paper Generator & Plagiarism Advisor API Endpoints ---

class PaperGenerateRequest(BaseModel):
    topic: str
    journal_tier: Optional[str] = "q1_ieee"
    journal_format: Optional[str] = "ieee"
    author_name: Optional[str] = "Manjunath"
    author_affiliation: Optional[str] = "Department of Artificial Intelligence and Data Science, JKK Munirajah College of Technology (JKKMCT), Tamil Nadu, India"

@app.get("/api/v1/generator/tiers")
def get_journal_tiers():
    return JOURNAL_TIERS

@app.post("/api/v1/generator/generate")
async def generate_paper(
    topic: str = Form(...),
    journal_tier: str = Form("q1_ieee"),
    journal_format: str = Form("ieee"),
    author_name: str = Form("Manjunath"),
    author_affiliation: str = Form("Department of Artificial Intelligence and Data Science, JKK Munirajah College of Technology (JKKMCT), Tamil Nadu, India"),
    notes_file: Optional[UploadFile] = File(None)
):
    if not topic.strip():
        raise HTTPException(status_code=400, detail="Topic string cannot be empty.")
    
    context_notes = None
    if notes_file:
        ext = os.path.splitext(notes_file.filename)[1].lower()
        temp_id = str(uuid.uuid4())
        temp_path = os.path.join(settings.UPLOAD_DIR, f"notes_{temp_id}{ext}")
        
        with open(temp_path, "wb") as buffer:
            shutil.copyfileobj(notes_file.file, buffer)
            
        try:
            if ext == ".txt":
                with open(temp_path, "r", encoding="utf-8", errors="ignore") as f:
                    context_notes = f.read()
            elif ext == ".pdf":
                from backend.app.parser.pdf_parser import parse_pdf
                parsed = parse_pdf(temp_path)
                context_notes = "\n".join([block.get("original_text", "") for block in parsed])
            elif ext == ".docx":
                from backend.app.parser.docx_parser import parse_docx
                parsed = parse_docx(temp_path)
                context_notes = "\n".join([block.get("original_text", "") for block in parsed])
            elif ext in [".pptx", ".ppt"]:
                try:
                    from pptx import Presentation
                    prs = Presentation(temp_path)
                    text_runs = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text_runs.append(shape.text)
                    context_notes = "\n".join(text_runs)
                except Exception as e:
                    print(f"Error parsing PPTX: {e}")
                    context_notes = f"[Parsed slides summary for {topic}]"
        except Exception as e:
            print(f"Error reading supplementary file: {e}")
        finally:
            if os.path.exists(temp_path):
                os.remove(temp_path)

    return generate_full_paper(
        topic=topic,
        journal_tier=journal_tier or "q1_ieee",
        journal_format=journal_format or "ieee",
        author_name=author_name or "Manjunath",
        author_affiliation=author_affiliation or "Department of Artificial Intelligence and Data Science, JKK Munirajah College of Technology (JKKMCT), Tamil Nadu, India",
        context_notes=context_notes
    )

from backend.app.advisor.research_agent import research_agent

@app.get("/api/v1/advisor/{paper_id}")
def get_plagiarism_advisor(paper_id: str):
    paper = db.get_paper(paper_id)
    if not paper:
        raise HTTPException(status_code=404, detail="Paper not found.")
    return generate_reduction_advice(paper_id)

# --- Autonomous Research Agent Endpoints ---
@app.get("/api/v1/research/search")
async def research_search(query: str = Query(...), limit: int = Query(6)):
    if not query.strip():
        raise HTTPException(status_code=400, detail="Query parameter cannot be empty.")
    results = await research_agent.search_literature(query=query, limit=limit)
    return {"query": query, "count": len(results), "papers": results}

@app.post("/api/v1/research/synthesize")
async def research_synthesize(topic: str = Form(...), style: str = Form("numeric")):
    if not topic.strip():
        raise HTTPException(status_code=400, detail="Topic cannot be empty.")
    return await research_agent.synthesize_literature_review(topic=topic, style=style)

@app.post("/api/v1/research/grounding")
async def research_grounding(passage: str = Form(...)):
    if not passage.strip():
        raise HTTPException(status_code=400, detail="Passage text cannot be empty.")
    return await research_agent.verify_claim_grounding(passage=passage)



